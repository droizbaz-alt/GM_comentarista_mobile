import os
import subprocess
import sys

# Asegurar python-pptx
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
except ImportError:
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx", "--quiet"])
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

def set_dark_bg(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(17, 24, 39) # Dark slate/navy

def create_pro_presentation():
    # Rutas absolutas a las imágenes generadas por IA
    img_brain = r"C:\Users\droiz\.gemini\antigravity\brain\841a34ef-f8ce-4b08-86bd-f0a44f1d1681\ia_chess_brain_1773915261296.png"
    img_mobile = r"C:\Users\droiz\.gemini\antigravity\brain\841a34ef-f8ce-4b08-86bd-f0a44f1d1681\mobile_chess_app_1773915275142.png"
    img_board = r"C:\Users\droiz\.gemini\antigravity\brain\841a34ef-f8ce-4b08-86bd-f0a44f1d1681\chess_board_abstract_1773915291894.png"

    prs = Presentation()
    # 16:9 ratio
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6] # Blank slide

    def add_slide_split(title_text, body_paragraphs, img_path=None, reverse=False):
        slide = prs.slides.add_slide(blank_layout)
        set_dark_bg(slide)

        # Dimensiones de la división
        half_width = Inches(13.333 / 2)
        full_height = Inches(7.5)

        # Renderizar imagen
        if img_path and os.path.exists(img_path):
            img_x = 0 if reverse else half_width
            slide.shapes.add_picture(img_path, img_x, 0, width=half_width, height=full_height)

        # Título
        txBox = slide.shapes.add_textbox(Inches(0.5) if not reverse else half_width + Inches(0.5), Inches(1), half_width - Inches(1), Inches(1))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.bold = True
        p.font.size = Pt(40)
        p.font.color.rgb = RGBColor(255, 255, 255) # White
        p.font.name = "Segoe UI"
        
        # Subrayado decorativo
        line = slide.shapes.add_shape(
            1, # Line
            Inches(0.5) if not reverse else half_width + Inches(0.5), Inches(2), Inches(1), Inches(0))
        line.line.color.rgb = RGBColor(74, 144, 217) # Azul accent
        line.line.width = Pt(4)

        # Cuerpo
        bodyBox = slide.shapes.add_textbox(Inches(0.5) if not reverse else half_width + Inches(0.5), Inches(2.5), half_width - Inches(1), Inches(4))
        bf = bodyBox.text_frame
        bf.word_wrap = True
        
        for i, text in enumerate(body_paragraphs):
            bp = bf.add_paragraph() if i > 0 else bf.paragraphs[0]
            bp.text = text
            bp.font.size = Pt(22)
            bp.font.color.rgb = RGBColor(209, 213, 219) # Gris claro
            bp.font.name = "Segoe UI"
            bp.space_after = Pt(14)

        return slide

    def add_slide_full_center(title_text, subtitle_text, img_bg=None):
        slide = prs.slides.add_slide(blank_layout)
        
        if img_bg and os.path.exists(img_bg):
            slide.shapes.add_picture(img_bg, 0, 0, prs.slide_width, prs.slide_height)
            # Para oscuro de fondo en lugar de transparente, ponemos un rectángulo negro semitransparente... 
            # Como pptx no soporta opacity del fill fácil, pondremos el texto en una franja oscura.
            strip = slide.shapes.add_shape(1, 0, Inches(2.5), prs.slide_width, Inches(2.5))
            strip.fill.solid()
            strip.fill.fore_color.rgb = RGBColor(0, 0, 0)
            strip.line.fill.background()
        else:
            set_dark_bg(slide)

        txBox = slide.shapes.add_textbox(0, Inches(2.8), prs.slide_width, Inches(1))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        p.alignment = PP_ALIGN.CENTER
        p.font.bold = True
        p.font.size = Pt(54)
        p.font.color.rgb = RGBColor(255, 255, 255) 

        subBox = slide.shapes.add_textbox(0, Inches(4), prs.slide_width, Inches(1))
        sf = subBox.text_frame
        sp = sf.paragraphs[0]
        sp.text = subtitle_text
        sp.alignment = PP_ALIGN.CENTER
        sp.font.size = Pt(28)
        sp.font.color.rgb = RGBColor(74, 144, 217)

    # PORTADA
    add_slide_full_center("GM Comentarista Mobile", "Tu Gran Maestro de bolsillo impulsado por Inteligencia Artificial", img_brain)

    # 2: Problema
    add_slide_split(
        "El muro del aprendizaje", 
        [
            "• Los motores (Stockfish) dan números fríos: +2.5, -4.0",
            "• Los módulos te dicen CUÁL es la mejor jugada, pero rara vez el POR QUÉ.",
            "• Los jugadores no entienden su propio error solo viendo una flecha."
        ], 
        img_board, 
        reverse=False
    )

    # 3: Solución
    add_slide_split(
        "¿Qué es GM Comentarista?", 
        [
            "✓ App Progresiva (PWA): Funciona en la nube sin agotar tu batería.",
            "✓ Análisis Dual: Stockfish evalúa, Gemini traduce.",
            "✓ Explicación Pedagógica: La app actúa como tu entrenador personal en tu mismo idioma."
        ],
        img_mobile,
        reverse=True
    )

    # 4: Uso
    add_slide_split(
        "3 Pasos Rápidos", 
        [
            "1. Importación Fácil: Conecta a Lichess o sube un PGN.",
            "2. Configuración Adaptable: Elige rigor (Rápido, Estándar, Profundo).",
            "3. Lectura Guiada: Un visor interactivo te lleva por tus errores."
        ],
        img_brain, # Reutilizamos imagen de la IA
        reverse=False
    )

    # 5: Tech
    add_slide_split(
        "Arquitectura & IA", 
        [
            "• Frontend: Streamlit ('Mobile First').",
            "• IA Híbrida: Modelo Lite para jugadas normales, Modelo Pro para blunders.",
            "• Rendimiento: Caché inteligente para minimizar uso de API."
        ],
        img_board,
        reverse=True
    )

    # 6: Cierre
    add_slide_full_center("¿Listos para subir de ELO?", "Exporta tu PGN guiado a ChessBase o Lichess")

    output_path = "PRESENTACION_PROFESIONAL.pptx"
    prs.save(output_path)
    print(f"Presentación PRO guardada en: {output_path}")

if __name__ == "__main__":
    create_pro_presentation()
