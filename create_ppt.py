import subprocess
import sys

# Instalar python-pptx si no existe
try:
    from pptx import Presentation
except ImportError:
    print("Instalando python-pptx...")
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx", "--quiet"])
    from pptx import Presentation

def create_presentation():
    prs = Presentation()

    # 1: Portada
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "GM Comentarista Mobile"
    slide.placeholders[1].text = "Tu Gran Maestro de bolsillo impulsado por Inteligencia Artificial"

    # 2: Problema
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "El muro del aprendizaje en ajedrez"
    tf = slide.placeholders[1].text_frame
    tf.text = "Los motores (Stockfish) dan números fríos: +2.5, -4.0"
    tf.add_paragraph().text = "Los módulos te dicen *cuál* es la mejor jugada, pero rara vez el *por qué*."
    tf.add_paragraph().text = "Los entrenadores humanos son excelentes, pero no están disponibles 24/7."
    tf.add_paragraph().text = "Los jugadores casuales no entienden su propio error solo viendo una flecha."

    # 3: Solución
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "¿Qué es GM Comentarista?"
    tf = slide.placeholders[1].text_frame
    tf.text = "App Progresiva (PWA): Funciona en la nube sin agotar la batería del teléfono."
    tf.add_paragraph().text = "Análisis Dual:\n  1. Stockfish evalúa matemáticamente.\n  2. Modelos de Lenguaje (Gemini/Gemma) lo traducen a lenguaje natural."
    tf.add_paragraph().text = "Explicación Pedagógica: La app actúa como un entrenador en tu idioma, adaptándose a la gravedad de tu error."

    # 4: Uso
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "¿Cómo funciona en 3 simples pasos?"
    tf = slide.placeholders[1].text_frame
    tf.text = "1. Importación Fácil: Conecta tu usuario de Lichess o sube un PGN clásico con un toque."
    tf.add_paragraph().text = "2. Configuración Adaptable: Elige el rigor del análisis (Rápido, Estándar o Profundo)."
    tf.add_paragraph().text = "3. Lectura Guiada: El visor interactivo te lleva de la mano repasando y explicando tus errores."

    # 5: Tech
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Bajo el capó: Arquitectura y Tecnología"
    tf = slide.placeholders[1].text_frame
    tf.text = "Frontend Mobile-First: Construido en Streamlit con diseño responsive."
    tf.add_paragraph().text = "Motor Lógico Inyectado: Binarios optimizados de Stockfish subyacentes."
    tf.add_paragraph().text = "IA Híbrida: \n  • Modelo 'Lite' (muy veloz) para el avance de la partida normal.\n  • Modelo 'Pro' (inteligente) reservado para explicar los errores graves de manera elocuente."
    tf.add_paragraph().text = "Caché de Posiciones: Minimiza costes evitando cálculos de API redundantes."

    # 6: Exportación
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Tu conocimiento, donde tú quieras"
    tf = slide.placeholders[1].text_frame
    tf.text = "El análisis no se queda encerrado dentro de la aplicación."
    tf.add_paragraph().text = "Inyección PGN: La IA inserta las explicaciones en los campos estándar de comentarios PGN."
    tf.add_paragraph().text = "Libertad total: Descarga el archivo PGN y léelo después en ChessBase o en Lichess."

    # 7: Cierre
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "¿Listos para subir de ELO?"
    tf = slide.placeholders[1].text_frame
    tf.text = "Beneficio principal: Optimización drástica en el tiempo de estudio."
    tf.add_paragraph().text = "Comprensión profunda: Memorizar variantes es menos útil que entender los planes subyacentes."
    tf.add_paragraph().text = "Q&A y Demo en vivo del análisis de una posición crítica."

    output_path = "PRESENTACION_DIAPOSITIVAS.pptx"
    prs.save(output_path)
    print(f"Presentación guardada exitosamente en: {output_path}")

if __name__ == "__main__":
    create_presentation()
